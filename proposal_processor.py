#!/usr/bin/env python3
"""
PowerPoint Proposal Generator for TCNG Image Editor
Handles image placeholder {{TP_MSB}} replacement and form data processing
"""

import argparse
import json
import os
import sys
from datetime import datetime
import tempfile

# Fix for Python 3.12+ compatibility with python-pptx
try:
    import collections.abc
    import collections
    if not hasattr(collections, 'Container'):
        collections.Container = collections.abc.Container
    if not hasattr(collections, 'Iterable'):
        collections.Iterable = collections.abc.Iterable  
    if not hasattr(collections, 'Mapping'):
        collections.Mapping = collections.abc.Mapping
    if not hasattr(collections, 'MutableMapping'):
        collections.MutableMapping = collections.abc.MutableMapping
    if not hasattr(collections, 'Sequence'):
        collections.Sequence = collections.abc.Sequence
except ImportError:
    pass

# Now import python-pptx
from pptx import Presentation
from pptx.util import Cm
from PIL import Image

def resize_image_to_powerpoint_dimensions(image_path, width_cm, height_cm, suffix=''):
    """
    Resize image to exact PowerPoint dimensions
    Args:
        image_path: Path to the input image
        width_cm: Target width in centimeters
        height_cm: Target height in centimeters
        suffix: Optional suffix for temporary file naming
    Returns:
        Path to resized image
    """
    try:
        # Open and process image
        with Image.open(image_path) as img:
            # Convert to RGB if necessary
            if img.mode not in ('RGB', 'RGBA'):
                img = img.convert('RGB')
            
            # Calculate target pixel dimensions (300 DPI for high quality)
            dpi = 300
            target_width_px = int(width_cm * dpi / 2.54)
            target_height_px = int(height_cm * dpi / 2.54)
            
            # Resize image
            resized_img = img.resize((target_width_px, target_height_px), Image.Resampling.LANCZOS)
            
            # Save to temporary file
            temp_fd, temp_path = tempfile.mkstemp(suffix='.png', prefix=f'resized_{suffix}_')
            os.close(temp_fd)  # Close file descriptor, we'll use the path
            
            resized_img.save(temp_path, 'PNG', optimize=True, dpi=(dpi, dpi))
            
            print(f"✅ Image resized to {target_width_px}x{target_height_px}px ({width_cm}x{height_cm}cm)")
            return temp_path
            
    except Exception as e:
        print(f"❌ Error resizing image: {str(e)}")
        return None

def replace_placeholders_in_pptx(template_path, form_data, msb_image_path, mccb_image_path, output_path):
    """
    Replace placeholders in PowerPoint template and insert images
    Args:
        template_path: Path to PowerPoint template
        form_data: Dictionary containing form data
        msb_image_path: Path to the MSB image file (can be None)
        mccb_image_path: Path to the MCCB image file (can be None)
        output_path: Path where to save the output file
    """
    try:
        print("📖 Loading PowerPoint template...")
        print(f"📁 Template path: {template_path}")
        print(f"🖼️ MSB Image path: {msb_image_path}")
        print(f"🖼️ MCCB Image path: {mccb_image_path}")
        print(f"📊 Form data: {form_data}")
        
        prs = Presentation(template_path)
        print(f"✅ Loaded presentation with {len(prs.slides)} slides")
        
        # First, let's inspect all slides and shapes to find what we're working with
        print("\n🔍 DEBUGGING: Inspecting all slides and shapes...")
        for slide_num, slide in enumerate(prs.slides):
            print(f"\n--- SLIDE {slide_num + 1} ---")
            print(f"Number of shapes: {len(slide.shapes)}")
            
            for shape_idx, shape in enumerate(slide.shapes):
                print(f"  Shape {shape_idx + 1}:")
                print(f"    Type: {shape.shape_type}")
                print(f"    Has name attr: {hasattr(shape, 'name')}")
                
                if hasattr(shape, 'name'):
                    print(f"    Name: '{shape.name}'")
                    if shape.name == '{{TP_MSB}}' or shape.name == 'TP_MSB':
                        print(f"    *** FOUND {{TP_MSB}} PLACEHOLDER! ***")
                        print(f"    Position: left={shape.left}, top={shape.top}")
                        print(f"    Size: width={shape.width}, height={shape.height}")
                    elif shape.name == '{{TP_MCCB}}' or shape.name == 'TP_MCCB':
                        print(f"    *** FOUND {{TP_MCCB}} PLACEHOLDER! ***")
                        print(f"    Position: left={shape.left}, top={shape.top}")
                        print(f"    Size: width={shape.width}, height={shape.height}")
                
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    print(f"    Has text_frame: True")
                    all_text = ""
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            all_text += run.text
                    if all_text.strip():
                        print(f"    Text content: '{all_text[:100]}{'...' if len(all_text) > 100 else ''}'")
                        if '{{TP_MSB}}' in all_text or 'TP_MSB' in all_text:
                            print(f"    *** CONTAINS {{TP_MSB}} TEXT! ***")
                        elif '{{TP_MCCB}}' in all_text or 'TP_MCCB' in all_text:
                            print(f"    *** CONTAINS {{TP_MCCB}} TEXT! ***")
                
                if hasattr(shape, 'text') and shape.text:
                    print(f"    Direct text: '{shape.text[:100]}{'...' if len(shape.text) > 100 else ''}'")
                    if '{{TP_MSB}}' in shape.text or 'TP_MSB' in shape.text:
                        print(f"    *** CONTAINS {{TP_MSB}} TEXT! ***")
                    elif '{{TP_MCCB}}' in shape.text or 'TP_MCCB' in shape.text:
                        print(f"    *** CONTAINS {{TP_MCCB}} TEXT! ***")
        
        # Prepare replacement mappings
        replacements = {
            '{{BUILDINGNAME}}': form_data.get('building_name', ''),
            '{{ADDRESS}}': form_data.get('address', ''),
        }
        
        print(f"\n🔄 Processing image replacements...")
        
        # Define image placeholders and their dimensions
        image_placeholders = [
            {
                'placeholder': '{{TP_MSB}}',
                'alt_placeholder': 'TP_MSB', 
                'image_path': msb_image_path,
                'width_cm': 9.05,
                'height_cm': 9.25,
                'suffix': 'msb'
            },
            {
                'placeholder': '{{TP_MCCB}}',
                'alt_placeholder': 'TP_MCCB',
                'image_path': mccb_image_path, 
                'width_cm': 4.22,
                'height_cm': 4.57,
                'suffix': 'mccb'
            }
        ]
        
        total_replacements_made = 0
        
        # Process each image placeholder
        for placeholder_info in image_placeholders:
            if placeholder_info['image_path'] and os.path.exists(placeholder_info['image_path']):
                print(f"🖼️ Processing {placeholder_info['placeholder']} image: {placeholder_info['image_path']}")
                print(f"📏 Image file size: {os.path.getsize(placeholder_info['image_path'])} bytes")
                
                # Resize image to exact PowerPoint dimensions
                resized_image_path = resize_image_to_powerpoint_dimensions(
                    placeholder_info['image_path'], 
                    width_cm=placeholder_info['width_cm'],
                    height_cm=placeholder_info['height_cm'],
                    suffix=placeholder_info['suffix']
                )
                
                if resized_image_path:
                    print(f"✅ Resized image created: {resized_image_path}")
                    print(f"📏 Resized image file size: {os.path.getsize(resized_image_path)} bytes")
                    
                    replacements_made = 0
                    
                    # Process slides to find and replace image placeholder
                    for slide_num, slide in enumerate(prs.slides):
                        print(f"\n🔍 Processing slide {slide_num + 1} for {placeholder_info['placeholder']}...")
                        shapes_to_remove = []
                        image_positions = []
                        
                        # Find placeholder by name
                        print(f"  Looking for shapes named '{placeholder_info['placeholder']}'...")
                        for shape_idx, shape in enumerate(slide.shapes):
                            if hasattr(shape, 'name'):
                                print(f"    Shape {shape_idx + 1} name: '{shape.name}'")
                                if (shape.name == placeholder_info['placeholder'] or 
                                    shape.name == placeholder_info['alt_placeholder']):
                                    print(f"🎯 Found {placeholder_info['placeholder']} placeholder on slide {slide_num + 1}")
                                    print(f"    Position: left={shape.left}, top={shape.top}")
                                    print(f"    Size: width={shape.width}, height={shape.height}")
                                    image_positions.append({
                                        'left': shape.left,
                                        'top': shape.top,
                                        'width': shape.width,
                                        'height': shape.height,
                                        'slide': slide
                                    })
                                    shapes_to_remove.append(shape)
                    
                        # Also check for text placeholders
                        print(f"  Looking for text containing '{placeholder_info['placeholder']}'...")
                        for shape_idx, shape in enumerate(slide.shapes):
                            shape_has_placeholder = False
                            
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if (placeholder_info['placeholder'] in run.text or 
                                            placeholder_info['alt_placeholder'] in run.text):
                                            print(f"🎯 Found {placeholder_info['placeholder']} text placeholder on slide {slide_num + 1}")
                                            print(f"    Shape type: {shape.shape_type}")
                                            print(f"    Full text: '{run.text}'")
                                            print(f"    Shape position: left={shape.left}, top={shape.top}")
                                            print(f"    Shape size: width={shape.width}, height={shape.height}")
                                            
                                            # If text contains only the placeholder, treat this shape as image placeholder
                                            if (run.text.strip() == placeholder_info['placeholder'] or 
                                                run.text.strip() == placeholder_info['alt_placeholder']):
                                                print("    → This appears to be a text-based image placeholder!")
                                                image_positions.append({
                                                    'left': shape.left,
                                                    'top': shape.top,
                                                    'width': shape.width,
                                                    'height': shape.height,
                                                    'slide': slide
                                                })
                                                shapes_to_remove.append(shape)
                                                shape_has_placeholder = True
                                            else:
                                                # Remove the text placeholder but keep the shape
                                                run.text = run.text.replace(placeholder_info['placeholder'], '')
                                                run.text = run.text.replace(placeholder_info['alt_placeholder'], '')
                                                print(f"    Text after replacement: '{run.text}'")
                            
                            if (not shape_has_placeholder and hasattr(shape, 'text') and 
                                (placeholder_info['placeholder'] in shape.text or 
                                 placeholder_info['alt_placeholder'] in shape.text)):
                                print(f"🎯 Found {placeholder_info['placeholder']} in direct text on slide {slide_num + 1}")
                                print(f"    Shape type: {shape.shape_type}")
                                print(f"    Full text: '{shape.text}'")
                                print(f"    Shape position: left={shape.left}, top={shape.top}")
                                print(f"    Shape size: width={shape.width}, height={shape.height}")
                                
                                # If text contains only the placeholder, treat this shape as image placeholder
                                if (shape.text.strip() == placeholder_info['placeholder'] or 
                                    shape.text.strip() == placeholder_info['alt_placeholder']):
                                    print("    → This appears to be a text-based image placeholder!")
                                    image_positions.append({
                                        'left': shape.left,
                                        'top': shape.top,
                                        'width': shape.width,
                                        'height': shape.height,
                                        'slide': slide
                                    })
                                    shapes_to_remove.append(shape)
                                else:
                                    shape.text = shape.text.replace(placeholder_info['placeholder'], '')
                                    shape.text = shape.text.replace(placeholder_info['alt_placeholder'], '')
                                    print(f"    Text after replacement: '{shape.text}'")
                    
                        
                        # Remove placeholder shapes
                        print(f"  Removing {len(shapes_to_remove)} placeholder shapes...")
                        for shape in shapes_to_remove:
                            try:
                                slide.shapes._spTree.remove(shape._element)
                                print("  ✅ Removed placeholder shape")
                            except Exception as e:
                                print(f"  ⚠️ Could not remove placeholder shape: {e}")
                        
                        # Add images at the stored positions
                        print(f"  Adding images at {len(image_positions)} positions...")
                        for pos_idx, pos_info in enumerate(image_positions):
                            try:
                                print(f"    Position {pos_idx + 1}:")
                                print(f"      left={pos_info['left']}, top={pos_info['top']}")
                                print(f"      width={pos_info['width']}, height={pos_info['height']}")
                                
                                new_picture = pos_info['slide'].shapes.add_picture(
                                    resized_image_path,
                                    pos_info['left'],
                                    pos_info['top'],
                                    pos_info['width'],
                                    pos_info['height']
                                )
                                
                                # Bring image to front instead of sending to back
                                pic_element = new_picture._element
                                spTree = pos_info['slide'].shapes._spTree
                                # Remove from current position
                                spTree.remove(pic_element)
                                # Add to the end (front-most layer)
                                spTree.append(pic_element)
                                
                                print(f"  ✅ Inserted {placeholder_info['placeholder']} image on slide {slide_num + 1} at position {pos_idx + 1}")
                                replacements_made += 1
                                
                            except Exception as img_error:
                                print(f"  ❌ Could not insert {placeholder_info['placeholder']} image on slide {slide_num + 1}: {img_error}")
                                import traceback
                                traceback.print_exc()
                    
                    print(f"\n📊 Total {placeholder_info['placeholder']} replacements made: {replacements_made}")
                    total_replacements_made += replacements_made
                    
                    # Clean up resized image
                    try:
                        os.unlink(resized_image_path)
                        print("🗑️ Cleaned up temporary resized image")
                    except Exception as cleanup_error:
                        print(f"⚠️ Could not clean up resized image: {cleanup_error}")
                else:
                    print(f"❌ Failed to resize {placeholder_info['placeholder']} image, continuing without image")
            else:
                print(f"ℹ️ No {placeholder_info['placeholder']} image provided or image file not found")
        
        print(f"\n📊 Total image replacements made across all placeholders: {total_replacements_made}")
        
        # Process text replacements
        print(f"\n📝 Processing text replacements...")
        text_replacements_made = 0
        
        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                # Handle text frames (preserves formatting)
                if hasattr(shape, "text_frame"):
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            original_text = run.text
                            new_text = original_text
                            
                            # Replace placeholders
                            for placeholder, value in replacements.items():
                                if placeholder in new_text:
                                    new_text = new_text.replace(placeholder, value)
                                    print(f"📝 Replaced '{placeholder}' with '{value}' on slide {slide_num + 1}")
                                    text_replacements_made += 1
                            
                            if new_text != original_text:
                                run.text = new_text
        
        print(f"📊 Total text replacements made: {text_replacements_made}")
        
        # Save the presentation
        print(f"\n💾 Saving presentation to: {output_path}")
        prs.save(output_path)
        print("✅ Presentation saved successfully!")
        
        # Verify the output file was created
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"📁 Output file size: {file_size} bytes")
        else:
            print("❌ Output file was not created!")
        
        return True
        
    except Exception as e:
        print(f"❌ Error processing PowerPoint: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def format_date(date_string):
    """
    Format date string to readable format
    """
    if not date_string:
        return ''
    
    try:
        # Try parsing various date formats
        formats = ['%Y-%m-%d', '%Y/%m/%d', '%d-%m-%Y', '%d/%m/%Y']
        
        for fmt in formats:
            try:
                date_obj = datetime.strptime(date_string, fmt)
                return date_obj.strftime('%B %d, %Y')  # e.g., "January 15, 2024"
            except ValueError:
                continue
        
        # If no format worked, return original string
        return date_string
        
    except Exception:
        return date_string

def main():
    parser = argparse.ArgumentParser(description='Generate PowerPoint proposal from template and form data')
    parser.add_argument('--template', required=True, help='Path to PowerPoint template file')
    parser.add_argument('--output', required=True, help='Path for output PowerPoint file')
    parser.add_argument('--data', required=True, help='JSON string containing form data')
    parser.add_argument('--msb-image', help='Path to MSB image file (optional)', dest='msb_image')
    parser.add_argument('--mccb-image', help='Path to MCCB image file (optional)', dest='mccb_image')
    # Keep legacy --image argument for backward compatibility
    parser.add_argument('--image', help='Path to image file (legacy, maps to MSB image)', dest='legacy_image')
    
    args = parser.parse_args()
    
    try:
        # Parse form data
        form_data = json.loads(args.data)
        print("📋 Form data loaded successfully")
        
        # Validate template file
        if not os.path.exists(args.template):
            print(f"❌ Template file not found: {args.template}")
            sys.exit(1)
        
        # Handle legacy image argument (backward compatibility)
        msb_image_path = args.msb_image or args.legacy_image
        mccb_image_path = args.mccb_image
        
        # Validate image files if provided
        if msb_image_path and not os.path.exists(msb_image_path):
            print(f"⚠️ MSB image file not found: {msb_image_path}")
            msb_image_path = None
            
        if mccb_image_path and not os.path.exists(mccb_image_path):
            print(f"⚠️ MCCB image file not found: {mccb_image_path}")
            mccb_image_path = None
        
        # Create output directory if it doesn't exist
        output_dir = os.path.dirname(args.output)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)
        
        # Process the presentation
        success = replace_placeholders_in_pptx(
            args.template,
            form_data,
            msb_image_path,
            mccb_image_path,
            args.output
        )
        
        if success:
            print("🎉 Proposal generation completed successfully!")
            sys.exit(0)
        else:
            print("❌ Proposal generation failed!")
            sys.exit(1)
            
    except json.JSONDecodeError as e:
        print(f"❌ Invalid JSON data: {e}")
        sys.exit(1)
    except Exception as e:
        print(f"❌ Unexpected error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

if __name__ == '__main__':
    main()