# proposal_generator.py
# Complete PowerPoint Proposal Automation with Dual Image Editor

from flask import Flask, render_template, request, send_file, redirect, url_for, flash, jsonify
from pptx import Presentation
from pptx.util import Cm
import os
import webbrowser
import threading
import time
from datetime import datetime
import sys
from PIL import Image, ImageOps
import io
import base64
import uuid
import json

# Configuration
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "FTP_Template.pptx")
OUTPUT_FOLDER = "generated_proposals"
TEMP_IMAGES_FOLDER = "temp_images"

# Ensure folders exist
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
os.makedirs(TEMP_IMAGES_FOLDER, exist_ok=True)

app = Flask(__name__)
app.secret_key = 'your_secret_key_here'
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB max file size
app.config['UPLOAD_FOLDER'] = TEMP_IMAGES_FOLDER

# Image dimensions (matching your PowerPoint template)
# First image
TARGET_WIDTH_CM = 19.05
TARGET_HEIGHT_CM = 10.79

# Second image
TARGET_WIDTH_CM_2 = 17.69
TARGET_HEIGHT_CM_2 = 11.38

def replace_placeholders_and_images_in_pptx(template_path, form_data, image_path, image_path_2, output_path):
    """
    Replace placeholders and images in PowerPoint template while preserving formatting
    Now handles both IMG_PLACEHOLDER and IMG_PLACEHOLDER2
    """
    try:
        # Load the PowerPoint template
        prs = Presentation(template_path)
        
        # Text replacement mapping
        replacements = {
            '{{BUILDINGNAME}}': form_data.get('building_name', ''),
            '{{ADDRESS}}': form_data.get('address', ''),
            '{{SURVEYDATE}}': form_data.get('survey_date', ''),
            '{{PREPAREDBY}}': form_data.get('prepared_by', ''),
            '{{PREPAREDDATE}}': form_data.get('prepared_date', ''),
            '{{TYPEBUILDING}}': form_data.get('type_building', ''),
            '{{BUILDINGMANAGERNAME}}': form_data.get('building_manager_name', ''),
            '{{BUILDINGMANAGEREMAIL}}': form_data.get('building_manager_email', ''),
            '{{BUILDINGMANAGERPHONE}}': form_data.get('building_manager_phone', ''),
            '{{BUILDINGMANAGERCOMPANY}}': form_data.get('building_manager_company', ''),
            '{{OTIC}}': form_data.get('otic', ''),
            '{{TAPNEWORSPARE}}': form_data.get('tap_new_or_spare', ''),
            '{{TAPPINGLOCATION}}': form_data.get('tapping_location', ''),
            '{{TAPPINGLOCATIONLEVEL}}': form_data.get('tapping_location_level', ''),
            '{{SITEASSESTMENTMCCB}}': form_data.get('site_assessment_mccb', ''),
            '{{TNBMETER}}': form_data.get('tnb_meter', ''),
            '{{TNBNA}}': form_data.get('tnb_na', ''),
            '{{PARKINGLOCATION}}': form_data.get('parking_location', ''),
            '{{NOOFCHARGERS}}': '2',  # Fixed to 2 as requested
            '{{EVCHARGERMODEL}}': form_data.get('ev_charger_model', ''),
            '{{NETWORKSTRENGTH}}': form_data.get('network_strength', '')
        }
        
        # Process all slides for text replacement and image insertion
        for slide_num, slide in enumerate(prs.slides):
            # First pass: Handle text replacement
            for shape in slide.shapes:
                # Handle text frames (preserves formatting)
                if hasattr(shape, "text_frame"):
                    text_frame = shape.text_frame
                    
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            original_text = run.text
                            new_text = original_text
                            
                            # Replace placeholders
                            for placeholder, value in replacements.items():
                                if placeholder in new_text:
                                    new_text = new_text.replace(placeholder, value)
                                    print(f"Replaced '{placeholder}' with '{value}' on slide {slide_num + 1}")
                            
                            if new_text != original_text:
                                run.text = new_text
                
                # Handle direct text (legacy support)
                elif hasattr(shape, "text"):
                    original_text = shape.text
                    new_text = original_text
                    
                    for placeholder, value in replacements.items():
                        if placeholder in new_text:
                            new_text = new_text.replace(placeholder, value)
                    
                    if new_text != original_text:
                        shape.text = new_text
            
            # Second pass: Handle image replacement (after text to avoid interfering with indexing)
            shapes_to_remove = []
            image_positions = []
            image_positions_2 = []
            
            for shape_idx, shape in enumerate(slide.shapes):
                # Handle first image placeholder
                if hasattr(shape, 'name') and shape.name == 'IMG_PLACEHOLDER' and image_path:
                    # Store the shape's position and properties
                    image_positions.append({
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height,
                        'slide': slide,
                        'shape_idx': shape_idx
                    })
                    shapes_to_remove.append(shape)
                
                # Handle second image placeholder
                elif hasattr(shape, 'name') and shape.name == 'IMG_PLACEHOLDER2' and image_path_2:
                    # Store the shape's position and properties
                    image_positions_2.append({
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height,
                        'slide': slide,
                        'shape_idx': shape_idx
                    })
                    shapes_to_remove.append(shape)
            
            # Remove placeholder shapes (in reverse order to maintain indexing)
            for shape in reversed(shapes_to_remove):
                slide.shapes._spTree.remove(shape._element)
            
            # Add first images at the stored positions
            for pos_info in image_positions:
                try:
                    # Add the new image
                    new_picture = pos_info['slide'].shapes.add_picture(
                        image_path, 
                        pos_info['left'], 
                        pos_info['top'], 
                        pos_info['width'], 
                        pos_info['height']
                    )
                    
                    # Send image to back (behind all other elements)
                    # Get the shape element
                    pic_element = new_picture._element
                    
                    # Move to the beginning of the shape tree (sends to back)
                    spTree = pos_info['slide'].shapes._spTree
                    spTree.insert(2, pic_element)  # Position 2 is behind most content but after background
                    
                    print(f"✅ Inserted IMG_PLACEHOLDER on slide {slide_num + 1}")
                    
                except Exception as img_error:
                    print(f"❌ Could not insert IMG_PLACEHOLDER on slide {slide_num + 1}: {img_error}")
            
            # Add second images at the stored positions
            for pos_info in image_positions_2:
                try:
                    # Add the new image
                    new_picture = pos_info['slide'].shapes.add_picture(
                        image_path_2, 
                        pos_info['left'], 
                        pos_info['top'], 
                        pos_info['width'], 
                        pos_info['height']
                    )
                    
                    # Send image to back (behind all other elements)
                    # Get the shape element
                    pic_element = new_picture._element
                    
                    # Move to the beginning of the shape tree (sends to back)
                    spTree = pos_info['slide'].shapes._spTree
                    spTree.insert(2, pic_element)  # Position 2 is behind most content but after background
                    
                    print(f"✅ Inserted IMG_PLACEHOLDER2 on slide {slide_num + 1}")
                    
                except Exception as img_error:
                    print(f"❌ Could not insert IMG_PLACEHOLDER2 on slide {slide_num + 1}: {img_error}")
        
        # Process tables for text replacement
        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    original_text = run.text
                                    new_text = original_text
                                    
                                    for placeholder, value in replacements.items():
                                        if placeholder in new_text:
                                            new_text = new_text.replace(placeholder, value)
                                    
                                    if new_text != original_text:
                                        run.text = new_text
        
        # Save the customized presentation
        prs.save(output_path)
        return True, "Proposal generated successfully with preserved formatting and proper image layering!"
        
    except FileNotFoundError:
        return False, f"Template file not found: {template_path}"
    except Exception as e:
        return False, f"Error generating proposal: {str(e)}"

def process_cropped_image(image_data_url, crop_data, image_type="1"):
    """
    Process the cropped image from the web editor
    Now supports two different image types with different dimensions
    """
    try:
        # Decode base64 image
        if ',' in image_data_url:
            image_data = image_data_url.split(',')[1]
        else:
            image_data = image_data_url
            
        image_bytes = base64.b64decode(image_data)
        
        # Open image with PIL
        image = Image.open(io.BytesIO(image_bytes))
        
        # crop_data is already a dict from json.loads, no need to re-parse
        left = max(0, int(crop_data['x']))
        top = max(0, int(crop_data['y']))
        right = min(image.width, int(crop_data['x'] + crop_data['width']))
        bottom = min(image.height, int(crop_data['y'] + crop_data['height']))
        
        # Ensure we have a valid crop area
        if right <= left or bottom <= top:
            return None, "Invalid crop area - please adjust your selection"
        
        cropped_image = image.crop((left, top, right, bottom))
        
        # Resize to exact PowerPoint dimensions (convert cm to pixels at 300 DPI for quality)
        # Using higher DPI for better PowerPoint quality
        dpi = 300
        
        # Choose dimensions based on image type
        if image_type == "2":
            target_width_px = int(TARGET_WIDTH_CM_2 * dpi / 2.54)  # Convert cm to pixels
            target_height_px = int(TARGET_HEIGHT_CM_2 * dpi / 2.54)
        else:
            target_width_px = int(TARGET_WIDTH_CM * dpi / 2.54)  # Convert cm to pixels
            target_height_px = int(TARGET_HEIGHT_CM * dpi / 2.54)
        
        final_image = cropped_image.resize((target_width_px, target_height_px), Image.Resampling.LANCZOS)
        
        # Save processed image with high quality for PowerPoint
        temp_filename = f"processed_image_{image_type}_{uuid.uuid4().hex}.png"
        temp_path = os.path.join(TEMP_IMAGES_FOLDER, temp_filename)
        final_image.save(temp_path, 'PNG', optimize=True)
        
        print(f"✅ Image {image_type} processed: {target_width_px}×{target_height_px}px saved to {temp_filename}")
        
        return temp_path, None
        
    except Exception as e:
        print(f"❌ Image {image_type} processing error: {str(e)}")
        return None, f"Error processing image {image_type}: {str(e)}"

@app.route('/')
def index():
    """
    Main page with enhanced form including dual image editor - Modern #F454CD Theme
    """
    return '''
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Proposal Generator - Modern Interface</title>
    <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap" rel="stylesheet">
    <style>
        :root {
            --primary: #F454CD;
            --primary-dark: #D63CAD;
            --primary-light: #F97DE0;
            --primary-ultra-light: #FDF2FB;
            --secondary: #6366F1;
            --accent: #10B981;
            --neutral-50: #FAFAFA;
            --neutral-100: #F5F5F5;
            --neutral-200: #E5E5E5;
            --neutral-300: #D4D4D4;
            --neutral-400: #A3A3A3;
            --neutral-500: #737373;
            --neutral-600: #525252;
            --neutral-700: #404040;
            --neutral-800: #262626;
            --neutral-900: #171717;
            --success: #10B981;
            --warning: #F59E0B;
            --error: #EF4444;
            --shadow-sm: 0 1px 2px 0 rgb(0 0 0 / 0.05);
            --shadow-md: 0 4px 6px -1px rgb(0 0 0 / 0.1), 0 2px 4px -2px rgb(0 0 0 / 0.1);
            --shadow-lg: 0 10px 15px -3px rgb(0 0 0 / 0.1), 0 4px 6px -4px rgb(0 0 0 / 0.1);
            --shadow-xl: 0 20px 25px -5px rgb(0 0 0 / 0.1), 0 8px 10px -6px rgb(0 0 0 / 0.1);
            --radius-sm: 0.375rem;
            --radius-md: 0.5rem;
            --radius-lg: 0.75rem;
            --radius-xl: 1rem;
        }

        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }

        body {
            font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif;
            background: linear-gradient(135deg, var(--neutral-50) 0%, #FFFFFF 50%, var(--primary-ultra-light) 100%);
            min-height: 100vh;
            color: var(--neutral-700);
            line-height: 1.6;
        }

        .container {
            max-width: 1200px;
            margin: 0 auto;
            padding: 2rem;
        }

        /* Header */
        .header {
            text-align: center;
            margin-bottom: 3rem;
            padding: 2rem 0;
        }

        .header h1 {
            font-size: 3rem;
            font-weight: 700;
            background: linear-gradient(135deg, var(--primary) 0%, var(--primary-dark) 100%);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            background-clip: text;
            margin-bottom: 0.5rem;
            letter-spacing: -0.025em;
        }

        .header p {
            font-size: 1.125rem;
            color: var(--neutral-500);
            font-weight: 400;
        }

        /* Progress Steps */
        .progress-bar {
            display: flex;
            justify-content: space-between;
            margin-bottom: 3rem;
            position: relative;
        }

        .progress-bar::before {
            content: '';
            position: absolute;
            top: 1.5rem;
            left: 2rem;
            right: 2rem;
            height: 2px;
            background: var(--neutral-200);
            z-index: 1;
        }

        .step {
            display: flex;
            flex-direction: column;
            align-items: center;
            position: relative;
            z-index: 2;
            background: white;
            padding: 0 1rem;
        }

        .step-number {
            width: 3rem;
            height: 3rem;
            border-radius: 50%;
            background: var(--neutral-200);
            display: flex;
            align-items: center;
            justify-content: center;
            font-weight: 600;
            color: var(--neutral-500);
            margin-bottom: 0.5rem;
            transition: all 0.3s ease;
        }

        .step.active .step-number {
            background: var(--primary);
            color: white;
            transform: scale(1.1);
            box-shadow: 0 0 0 4px var(--primary-ultra-light);
        }

        .step-label {
            font-size: 0.875rem;
            font-weight: 500;
            color: var(--neutral-500);
            text-align: center;
        }

        .step.active .step-label {
            color: var(--primary);
        }

        /* Form Sections */
        .form-section {
            background: white;
            border-radius: var(--radius-xl);
            padding: 2rem;
            margin-bottom: 2rem;
            box-shadow: var(--shadow-sm);
            border: 1px solid var(--neutral-100);
            transition: all 0.3s ease;
        }

        .form-section:hover {
            box-shadow: var(--shadow-md);
            border-color: var(--primary-light);
        }

        .section-header {
            display: flex;
            align-items: center;
            margin-bottom: 1.5rem;
            padding-bottom: 1rem;
            border-bottom: 2px solid var(--neutral-100);
        }

        .section-icon {
            width: 2.5rem;
            height: 2.5rem;
            border-radius: var(--radius-md);
            background: linear-gradient(135deg, var(--primary) 0%, var(--primary-dark) 100%);
            display: flex;
            align-items: center;
            justify-content: center;
            margin-right: 1rem;
            font-size: 1.25rem;
        }

        .section-title {
            font-size: 1.5rem;
            font-weight: 600;
            color: var(--neutral-800);
            margin: 0;
        }

        /* Form Controls */
        .form-grid {
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(300px, 1fr));
            gap: 1.5rem;
        }

        .form-group {
            display: flex;
            flex-direction: column;
        }

        .form-label {
            font-size: 0.875rem;
            font-weight: 500;
            color: var(--neutral-700);
            margin-bottom: 0.5rem;
            display: flex;
            align-items: center;
        }

        .required {
            color: var(--error);
            margin-left: 0.25rem;
        }

        .form-input {
            padding: 0.875rem 1rem;
            border: 1.5px solid var(--neutral-200);
            border-radius: var(--radius-md);
            font-size: 0.875rem;
            transition: all 0.2s ease;
            background: white;
            font-family: inherit;
        }

        .form-input:focus {
            outline: none;
            border-color: var(--primary);
            box-shadow: 0 0 0 3px var(--primary-ultra-light);
        }

        .form-input::placeholder {
            color: var(--neutral-400);
        }

        .form-help {
            font-size: 0.75rem;
            color: var(--neutral-500);
            margin-top: 0.25rem;
        }

        /* Image Upload Sections */
        .dual-image-container {
            display: grid;
            grid-template-columns: 1fr;
            gap: 2rem;
            margin-top: 2rem;
        }

        .image-section {
            background: white;
            border-radius: var(--radius-xl);
            padding: 1.5rem;
            border: 2px dashed var(--neutral-200);
            transition: all 0.3s ease;
            position: relative;
            overflow: hidden;
        }

        .image-section.has-image {
            border-style: solid;
            border-color: var(--primary-light);
            background: var(--primary-ultra-light);
        }

        .image-section h3 {
            font-size: 1.125rem;
            font-weight: 600;
            color: var(--neutral-800);
            margin-bottom: 0.5rem;
            text-align: center;
        }

        .image-dimension {
            font-size: 0.75rem;
            color: var(--neutral-500);
            text-align: center;
            margin-bottom: 1rem;
            padding: 0.25rem 0.75rem;
            background: var(--neutral-100);
            border-radius: var(--radius-sm);
            display: inline-block;
            width: 100%;
        }

        .image-upload-area {
            border: 2px dashed var(--neutral-300);
            border-radius: var(--radius-lg);
            padding: 2rem 1rem;
            text-align: center;
            background: var(--neutral-50);
            transition: all 0.3s ease;
            cursor: pointer;
            position: relative;
        }

        .image-upload-area:hover {
            border-color: var(--primary);
            background: var(--primary-ultra-light);
        }

        .image-upload-area.dragover {
            border-color: var(--primary);
            background: var(--primary-ultra-light);
            transform: scale(1.02);
        }

        .upload-icon {
            font-size: 3rem;
            color: var(--neutral-400);
            margin-bottom: 1rem;
        }

        .upload-text {
            font-size: 1rem;
            font-weight: 500;
            color: var(--neutral-700);
            margin-bottom: 0.5rem;
        }

        .upload-subtext {
            font-size: 0.875rem;
            color: var(--neutral-500);
        }

        /* Image Editor */
        .image-editor {
            display: none;
            margin-top: 1rem;
            padding: 1rem;
            background: var(--neutral-50);
            border-radius: var(--radius-lg);
            border: 1px solid var(--neutral-200);
        }

        .editor-controls {
            display: flex;
            justify-content: space-between;
            align-items: center;
            margin-bottom: 1rem;
            padding: 1rem;
            background: white;
            border-radius: var(--radius-md);
            box-shadow: var(--shadow-sm);
        }

        .zoom-control {
            display: flex;
            align-items: center;
            gap: 0.75rem;
        }

        .zoom-label {
            font-size: 0.875rem;
            font-weight: 500;
            color: var(--neutral-600);
        }

        .zoom-slider {
            width: 150px;
            height: 6px;
            border-radius: 3px;
            background: var(--neutral-200);
            outline: none;
            -webkit-appearance: none;
        }

        .zoom-slider::-webkit-slider-thumb {
            -webkit-appearance: none;
            width: 20px;
            height: 20px;
            border-radius: 50%;
            background: var(--primary);
            cursor: pointer;
            box-shadow: var(--shadow-sm);
        }

        .zoom-slider::-moz-range-thumb {
            width: 20px;
            height: 20px;
            border-radius: 50%;
            background: var(--primary);
            cursor: pointer;
            border: none;
            box-shadow: var(--shadow-sm);
        }

        .zoom-value {
            font-size: 0.875rem;
            font-weight: 600;
            color: var(--primary);
            min-width: 3rem;
            text-align: center;
        }

        .editor-buttons {
            display: flex;
            gap: 0.5rem;
        }

        .btn {
            padding: 0.5rem 1rem;
            border: none;
            border-radius: var(--radius-md);
            font-size: 0.875rem;
            font-weight: 500;
            cursor: pointer;
            transition: all 0.2s ease;
            display: inline-flex;
            align-items: center;
            justify-content: center;
            gap: 0.5rem;
        }

        .btn-secondary {
            background: var(--neutral-100);
            color: var(--neutral-600);
            border: 1px solid var(--neutral-200);
        }

        .btn-secondary:hover {
            background: var(--neutral-200);
            color: var(--neutral-700);
        }

        .btn-primary {
            background: var(--primary);
            color: white;
        }

        .btn-primary:hover {
            background: var(--primary-dark);
            transform: translateY(-1px);
        }

        .canvas-container {
            position: relative;
            display: inline-block;
            background: white;
            border-radius: var(--radius-md);
            box-shadow: var(--shadow-md);
            overflow: hidden;
        }

        .image-canvas {
            cursor: grab;
            display: block;
        }

        .image-canvas:active {
            cursor: grabbing;
        }

        .debug-info {
            margin-top: 1rem;
            padding: 0.75rem;
            background: white;
            border-radius: var(--radius-md);
            border-left: 4px solid var(--primary);
            font-family: 'SF Mono', Monaco, monospace;
            font-size: 0.75rem;
            color: var(--neutral-600);
        }

        .editor-help {
            margin-top: 1rem;
            padding: 1rem;
            background: var(--primary-ultra-light);
            border-radius: var(--radius-md);
            border: 1px solid var(--primary-light);
        }

        .editor-help h4 {
            font-size: 0.875rem;
            font-weight: 600;
            color: var(--primary-dark);
            margin-bottom: 0.5rem;
        }

        .editor-help p {
            font-size: 0.8125rem;
            color: var(--neutral-600);
            line-height: 1.5;
        }

        /* Progress Upload */
        .upload-progress {
            display: none;
            margin-top: 1rem;
        }

        .progress-bar-container {
            background: var(--neutral-200);
            border-radius: var(--radius-sm);
            overflow: hidden;
            height: 8px;
            margin-bottom: 0.5rem;
        }

        .progress-bar-fill {
            background: linear-gradient(90deg, var(--primary) 0%, var(--primary-dark) 100%);
            height: 100%;
            width: 0%;
            transition: width 0.3s ease;
            border-radius: var(--radius-sm);
        }

        .progress-text {
            font-size: 0.875rem;
            color: var(--neutral-600);
            text-align: center;
        }

        /* Submit Button */
        .submit-section {
            background: white;
            border-radius: var(--radius-xl);
            padding: 2rem;
            box-shadow: var(--shadow-md);
            border: 1px solid var(--neutral-100);
            text-align: center;
            margin-top: 2rem;
        }

        .submit-btn {
            background: linear-gradient(135deg, var(--primary) 0%, var(--primary-dark) 100%);
            color: white;
            padding: 1rem 3rem;
            border: none;
            border-radius: var(--radius-lg);
            font-size: 1.125rem;
            font-weight: 600;
            cursor: pointer;
            transition: all 0.3s ease;
            box-shadow: var(--shadow-lg);
            position: relative;
            overflow: hidden;
            min-width: 250px;
        }

        .submit-btn::before {
            content: '';
            position: absolute;
            top: 0;
            left: -100%;
            width: 100%;
            height: 100%;
            background: linear-gradient(90deg, transparent, rgba(255,255,255,0.3), transparent);
            transition: left 0.5s ease;
        }

        .submit-btn:hover::before {
            left: 100%;
        }

        .submit-btn:hover {
            transform: translateY(-2px);
            box-shadow: var(--shadow-xl);
        }

        .submit-btn:disabled {
            background: var(--neutral-400);
            cursor: not-allowed;
            transform: none;
            box-shadow: var(--shadow-sm);
        }

        /* Animations */
        @keyframes slideInUp {
            from {
                opacity: 0;
                transform: translateY(30px);
            }
            to {
                opacity: 1;
                transform: translateY(0);
            }
        }

        .form-section {
            animation: slideInUp 0.6s ease forwards;
        }

        .form-section:nth-child(2) { animation-delay: 0.1s; }
        .form-section:nth-child(3) { animation-delay: 0.2s; }
        .form-section:nth-child(4) { animation-delay: 0.3s; }
        .form-section:nth-child(5) { animation-delay: 0.4s; }
        .form-section:nth-child(6) { animation-delay: 0.5s; }

        /* Mobile-Friendly Canvas */
        .canvas-container {
            position: relative;
            display: inline-block;
            background: white;
            border-radius: var(--radius-md);
            box-shadow: var(--shadow-md);
            overflow: hidden;
            touch-action: none; /* Prevent default touch behaviors */
            width: 100%;
            max-width: 100%;
        }

        .image-canvas {
            cursor: grab;
            display: block;
            max-width: 100%;
            height: auto;
            touch-action: none;
        }

        .image-canvas:active {
            cursor: grabbing;
        }

        /* Mobile-specific improvements */
        @media (max-width: 768px) {
            .container {
                padding: 0.5rem;
            }

            .header {
                padding: 1rem 0;
                margin-bottom: 2rem;
            }

            .header h1 {
                font-size: 2rem;
            }

            .form-section {
                padding: 1.5rem;
                margin-bottom: 1.5rem;
            }

            .form-grid {
                grid-template-columns: 1fr;
                gap: 1rem;
            }

            .progress-bar {
                flex-direction: column;
                gap: 1rem;
            }

            .progress-bar::before {
                display: none;
            }

            /* Mobile image editor improvements */
            .image-section {
                padding: 1rem;
            }

            .image-editor {
                padding: 0.5rem;
            }

            .editor-controls {
                flex-direction: column;
                gap: 1rem;
                align-items: stretch;
                padding: 0.75rem;
            }

            .zoom-control {
                justify-content: center;
                flex-wrap: wrap;
                gap: 0.5rem;
            }

            .zoom-slider {
                width: 200px;
                min-width: 150px;
            }

            .editor-buttons {
                justify-content: center;
                gap: 0.75rem;
            }

            .btn {
                padding: 0.75rem 1.5rem;
                font-size: 1rem;
                min-width: 100px;
            }

            /* Canvas responsive sizing */
            .canvas-container {
                width: 100%;
                overflow-x: auto;
            }

            .image-canvas {
                max-width: none;
                min-width: 300px;
            }

            /* Make form inputs more touch-friendly */
            .form-input {
                padding: 1rem;
                font-size: 1rem;
            }

            .submit-btn {
                padding: 1.25rem 2rem;
                font-size: 1.125rem;
                margin-top: 1rem;
            }

            /* Better spacing for mobile */
            .dual-image-container {
                gap: 1.5rem;
            }

            .section-header {
                flex-direction: column;
                align-items: flex-start;
                gap: 0.5rem;
                text-align: left;
            }

            .section-icon {
                margin-right: 0;
                margin-bottom: 0.5rem;
            }

            /* Mobile debug info */
            .debug-info {
                font-size: 0.7rem;
                padding: 0.5rem;
                overflow-x: auto;
            }

            .editor-help {
                padding: 0.75rem;
            }

            .editor-help h4 {
                font-size: 0.9rem;
            }

            .editor-help p {
                font-size: 0.8rem;
            }
        }

        /* Extra small mobile devices */
        @media (max-width: 480px) {
            .container {
                padding: 0.25rem;
            }

            .header h1 {
                font-size: 1.5rem;
            }

            .form-section {
                padding: 1rem;
            }

            .section-title {
                font-size: 1.25rem;
            }

            .canvas-container {
                margin: 0 auto;
            }

            .zoom-slider {
                width: 150px;
            }

            .editor-buttons {
                flex-direction: column;
            }

            .btn {
                width: 100%;
                margin-bottom: 0.5rem;
            }
        }

        /* Custom Scrollbar */
        ::-webkit-scrollbar {
            width: 8px;
        }

        ::-webkit-scrollbar-track {
            background: var(--neutral-100);
        }

        ::-webkit-scrollbar-thumb {
            background: var(--primary-light);
            border-radius: 4px;
        }

        ::-webkit-scrollbar-thumb:hover {
            background: var(--primary);
        }
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>Proposal Generator</h1>
            <p>Create professional proposals with custom images and detailed specifications</p>
        </div>

        <div class="progress-bar">
            <div class="step active">
                <div class="step-number">1</div>
                <div class="step-label">Project Info</div>
            </div>
            <div class="step">
                <div class="step-number">2</div>
                <div class="step-label">Manager Details</div>
            </div>
            <div class="step">
                <div class="step-number">3</div>
                <div class="step-label">Technical</div>
            </div>
            <div class="step">
                <div class="step-number">4</div>
                <div class="step-label">EV Chargers</div>
            </div>
            <div class="step">
                <div class="step-number">5</div>
                <div class="step-label">Images</div>
            </div>
        </div>
        
        <form id="proposalForm" method="POST" action="/generate" enctype="multipart/form-data">
            <!-- Project Information Section -->
            <div class="form-section">
                <div class="section-header">
                    <div class="section-icon">📋</div>
                    <h2 class="section-title">Project Information</h2>
                </div>
                
                <div class="form-grid">
                    <div class="form-group">
                        <label class="form-label" for="building_name">Building Name</label>
                        <input type="text" id="building_name" name="building_name" class="form-input" placeholder="Enter building name...">
                    </div>
                    
                    <div class="form-group">
                        <label class="form-label" for="type_building">Type of Building</label>
                        <select id="type_building" name="type_building" class="form-input">
                            <option value="">Select building type...</option>
                            <option value="Condominium">Condominium</option>
                            <option value="Office">Office</option>
                            <option value="Mall">Mall</option>
                            <option value="Hotel">Hotel</option>
                            <option value="Apartment">Apartment</option>
                            <option value="Commercial">Commercial</option>
                            <option value="Residential">Residential</option>
                            <option value="Mixed Development">Mixed Development</option>
                        </select>
                    </div>
                    
                    <div class="form-group" style="grid-column: 1 / -1;">
                        <label class="form-label" for="address">Project Address</label>
                        <textarea id="address" name="address" class="form-input" rows="2" placeholder="Enter the full project address..."></textarea>
                    </div>
                    
                    <div class="form-group">
                        <label class="form-label" for="survey_date">Survey Date</label>
                        <input type="date" id="survey_date" name="survey_date" class="form-input">
                    </div>
                    
                    <div class="form-group">
                        <label class="form-label" for="prepared_date">Prepared Date</label>
                        <input type="date" id="prepared_date" name="prepared_date" class="form-input">
                    </div>
                    
                    <div class="form-group">
                        <label class="form-label" for="prepared_by">Prepared By</label>
                        <input type="text" id="prepared_by" name="prepared_by" class="form-input" placeholder="Enter your name...">
                    </div>
                </div>
            </div>

            <!-- Building Manager Information Section -->
            <div class="form-section">
                <div class="section-header">
                    <div class="section-icon">👤</div>
                    <h2 class="section-title">Building Manager Information</h2>
                </div>
                
                <div class="form-grid">
                    <div class="form-group">
                        <label class="form-label" for="building_manager_name">Manager Name</label>
                        <input type="text" id="building_manager_name" name="building_manager_name" class="form-input" placeholder="Enter manager name...">
                    </div>
                    
                    <div class="form-group">
                        <label class="form-label" for="building_manager_company">Company</label>
                        <input type="text" id="building_manager_company" name="building_manager_company" class="form-input" placeholder="Enter company name...">
                    </div>
                    
                    <div class="form-group">
                        <label class="form-label" for="building_manager_email">Email Address</label>
                        <input type="email" id="building_manager_email" name="building_manager_email" class="form-input" placeholder="Enter email address...">
                    </div>
                    
                    <div class="form-group">
                        <label class="form-label" for="building_manager_phone">Phone Number</label>
                        <input type="text" id="building_manager_phone" name="building_manager_phone" class="form-input" placeholder="Enter phone number...">
                    </div>
                </div>
            </div>

            <!-- Technical Information Section -->
            <div class="form-section">
                <div class="section-header">
                    <div class="section-icon">⚡</div>
                    <h2 class="section-title">Technical Information</h2>
                </div>
                
                <div class="form-grid">
                    <div class="form-group">
                        <label class="form-label" for="otic">OTIC</label>
                        <input type="text" id="otic" name="otic" class="form-input" placeholder="Enter OTIC...">
                    </div>
                    
                    <div class="form-group">
                        <label class="form-label" for="tap_new_or_spare">Tap Type</label>
                        <select id="tap_new_or_spare" name="tap_new_or_spare" class="form-input">
                            <option value="Spare" selected>Spare</option>
                            <option value="Tap new">Tap new</option>
                            <option value="Replace">Replace</option>
                        </select>
                    </div>
                    
                    <div class="form-group">
                        <label class="form-label" for="tapping_location">Tapping Location</label>
                        <input type="text" id="tapping_location" name="tapping_location" class="form-input" placeholder="e.g. MSB, Feederpillar, SSB..." value="MSB">
                        <div class="form-help">Examples: MSB, Feederpillar, SSB</div>
                    </div>
                    
                    <div class="form-group">
                        <label class="form-label" for="tapping_location_level">Location Level</label>
                        <input type="text" id="tapping_location_level" name="tapping_location_level" class="form-input" placeholder="e.g. B1, L1, G...">
                        <div class="form-help">Examples: B1, L1, G, etc.</div>
                    </div>
                    
                    <div class="form-group" style="grid-column: 1 / -1;">
                        <label class="form-label" for="site_assessment_mccb">Site Assessment MCCB</label>
                        <input type="text" id="site_assessment_mccb" name="site_assessment_mccb" class="form-input" placeholder="Enter MCCB details..." value="1 x MCCB 100A TPN (New)">
                        <div class="form-help">Example: 1 x MCCB 100A TPN (New)</div>
                    </div>
                    
                    <div class="form-group">
                        <label class="form-label" for="tnb_meter">TNB Meter Type</label>
                        <select id="tnb_meter" name="tnb_meter" class="form-input" onchange="updateTnbNa()">
                            <option value="TNB Tenant Meter" selected>TNB Tenant Meter</option>
                            <option value="kWh Meter">kWh Meter</option>
                        </select>
                    </div>
                    
                    <div class="form-group">
                        <label class="form-label" for="tnb_na">TNB NA</label>
                        <input type="text" id="tnb_na" name="tnb_na" class="form-input" value="CNG" readonly style="background-color: var(--neutral-100);">
                        <div class="form-help">Automatically set based on TNB Meter Type</div>
                    </div>
                    
                    <div class="form-group">
                        <label class="form-label" for="network_strength">Network Strength</label>
                        <select id="network_strength" name="network_strength" class="form-input">
                            <option value="Average" selected>Average</option>
                            <option value="Excellent">Excellent</option>
                            <option value="Poor">Poor</option>
                            <option value="NA">NA</option>
                        </select>
                    </div>
                </div>
            </div>

            <!-- EV Charger Information Section -->
            <div class="form-section">
                <div class="section-header">
                    <div class="section-icon">🚗</div>
                    <h2 class="section-title">EV Charger Information</h2>
                </div>
                
                <div class="form-grid">
                    <div class="form-group">
                        <label class="form-label" for="ev_charger_model">Charger Model</label>
                        <select id="ev_charger_model" name="ev_charger_model" class="form-input">
                            <option value="22kW AC chargers" selected>22kW AC chargers</option>
                            <option value="7kW AC chargers">7kW AC chargers</option>
                            <option value="11kW AC chargers">11kW AC chargers</option>
                        </select>
                    </div>
                    
                    <div class="form-group">
                        <label class="form-label" for="no_of_chargers">Number of Chargers</label>
                        <input type="text" id="no_of_chargers" name="no_of_chargers" class="form-input" value="2" readonly style="background-color: var(--neutral-100);">
                        <div class="form-help">Fixed to 2 chargers for this version</div>
                    </div>
                    
                    <div class="form-group" style="grid-column: 1 / -1;">
                        <label class="form-label" for="parking_location">Parking Location</label>
                        <input type="text" id="parking_location" name="parking_location" class="form-input" placeholder="Enter lot no and location...">
                        <div class="form-help">Should include lot number and location details</div>
                    </div>
                </div>
            </div>
            
            <!-- Dual Image Upload and Editor Section -->
            <div class="form-section">
                <div class="section-header">
                    <div class="section-icon">🖼️</div>
                    <h2 class="section-title">Image Upload & Position Editor</h2>
                </div>
                <p style="margin-bottom: 2rem; color: var(--neutral-500);">Upload and position images for both placeholders in your template</p>
                
                <div class="dual-image-container">
                    <!-- First Image Section -->
                    <div class="image-section" id="imageSection1">
                        <h3>Front Page Building Photo</h3>
                        <div class="image-dimension">19.05 cm × 10.79 cm</div>
                        <p style="text-align: center; color: var(--neutral-500); font-size: 0.875rem; margin-bottom: 1rem;">For IMG_PLACEHOLDER</p>
                        
                        <div class="image-upload-area" id="uploadArea1">
                            <div class="upload-icon">📁</div>
                            <div class="upload-text">Click to Upload Image or Drag & Drop</div>
                            <div class="upload-subtext">Supports: JPG, PNG, GIF (Max 50MB)</div>
                            <div class="upload-progress" id="uploadProgress1">
                                <div class="progress-bar-container">
                                    <div class="progress-bar-fill" id="progressBar1"></div>
                                </div>
                                <div class="progress-text" id="progressText1">Processing image...</div>
                            </div>
                            <input type="file" id="imageUpload1" accept="image/*" style="display: none;">
                        </div>
                        
                        <div class="image-editor" id="imageEditor1">
                            <div class="editor-controls">
                                <div class="zoom-control">
                                    <div class="zoom-label">Zoom:</div>
                                    <input type="range" id="zoomSlider1" class="zoom-slider" min="0.1" max="3" step="0.1" value="1">
                                    <div class="zoom-value" id="zoomValue1">100%</div>
                                </div>
                                <div class="editor-buttons">
                                    <button type="button" class="btn btn-secondary" onclick="resetPosition(1)">Reset</button>
                                    <button type="button" class="btn btn-primary" onclick="autoFit(1)">Auto Fit</button>
                                </div>
                            </div>
                            
                            <div class="canvas-container">
                                <canvas id="imageCanvas1" class="image-canvas" width="760" height="430"></canvas>
                            </div>
                            
                            <div class="debug-info" id="debugInfo1">
                                <strong>Crop Coordinates:</strong> <span id="cropCoords1">Position image to see coordinates</span>
                            </div>
                            
                            <div class="editor-help">
                                <h4>How to Use:</h4>
                                <p>Drag the image to position it within the frame. The frame shows exactly what will be inserted into IMG_PLACEHOLDER.</p>
                            </div>
                        </div>
                        
                        <!-- Hidden inputs for first image -->
                        <input type="hidden" id="croppedImageData1" name="cropped_image_data">
                        <input type="hidden" id="cropCoordinates1" name="crop_coordinates">
                    </div>
                    
                    <!-- Second Image Section -->
                    <div class="image-section" id="imageSection2">
                        <h3>GPS Map Photo</h3>
                        <div class="image-dimension">17.69 cm × 11.38 cm</div>
                        <p style="text-align: center; color: var(--neutral-500); font-size: 0.875rem; margin-bottom: 1rem;">For IMG_PLACEHOLDER2</p>
                        
                        <div class="image-upload-area" id="uploadArea2">
                            <div class="upload-icon">📁</div>
                            <div class="upload-text">Click to Upload Image or Drag & Drop</div>
                            <div class="upload-subtext">Supports: JPG, PNG, GIF (Max 50MB)</div>
                            <div class="upload-progress" id="uploadProgress2">
                                <div class="progress-bar-container">
                                    <div class="progress-bar-fill" id="progressBar2"></div>
                                </div>
                                <div class="progress-text" id="progressText2">Processing image...</div>
                            </div>
                            <input type="file" id="imageUpload2" accept="image/*" style="display: none;">
                        </div>
                        
                        <div class="image-editor" id="imageEditor2">
                            <div class="editor-controls">
                                <div class="zoom-control">
                                    <div class="zoom-label">Zoom:</div>
                                    <input type="range" id="zoomSlider2" class="zoom-slider" min="0.1" max="3" step="0.1" value="1">
                                    <div class="zoom-value" id="zoomValue2">100%</div>
                                </div>
                                <div class="editor-buttons">
                                    <button type="button" class="btn btn-secondary" onclick="resetPosition(2)">Reset</button>
                                    <button type="button" class="btn btn-primary" onclick="autoFit(2)">Auto Fit</button>
                                </div>
                            </div>
                            
                            <div class="canvas-container">
                                <canvas id="imageCanvas2" class="image-canvas" width="708" height="456"></canvas>
                            </div>
                            
                            <div class="debug-info" id="debugInfo2">
                                <strong>Crop Coordinates:</strong> <span id="cropCoords2">Position image to see coordinates</span>
                            </div>
                            
                            <div class="editor-help">
                                <h4>How to Use:</h4>
                                <p>Drag the image to position it within the frame. The frame shows exactly what will be inserted into IMG_PLACEHOLDER2.</p>
                            </div>
                        </div>
                        
                        <!-- Hidden inputs for second image -->
                        <input type="hidden" id="croppedImageData2" name="cropped_image_data_2">
                        <input type="hidden" id="cropCoordinates2" name="crop_coordinates_2">
                    </div>
                </div>
            </div>
            
            <div class="submit-section">
                <button type="submit" class="submit-btn" id="generateBtn">
                    Generate Proposal Template
                </button>
            </div>
        </form>
    </div>
    
    <script>
        // All your existing JavaScript code - preserved exactly
        
        // Dual Image Editor JavaScript - Drag to Position System
        let editors = {
            1: {
                canvas: null, ctx: null, originalImage: null, currentImage: null,
                isDragging: false, lastMouseX: 0, lastMouseY: 0,
                zoom: 1, offsetX: 0, offsetY: 0,
                FRAME_WIDTH: 380, FRAME_HEIGHT: 216, TARGET_RATIO: 19.05 / 10.79
            },
            2: {
                canvas: null, ctx: null, originalImage: null, currentImage: null,
                isDragging: false, lastMouseX: 0, lastMouseY: 0,
                zoom: 1, offsetX: 0, offsetY: 0,
                FRAME_WIDTH: 354, FRAME_HEIGHT: 228, TARGET_RATIO: 17.69 / 11.38
            }
        };
        
        document.addEventListener('DOMContentLoaded', function() {
            // Set default dates
            document.getElementById('prepared_date').valueAsDate = new Date();
            
            // Initialize both canvases
            for (let i = 1; i <= 2; i++) {
                initializeEditor(i);
            }
            
            // Initialize step progress
            updateStepProgress();
        });
        
        function updateStepProgress() {
            // Simple step progress based on form completion
            // This could be enhanced to actually track form completion
            const steps = document.querySelectorAll('.step');
            steps.forEach((step, index) => {
                if (index === 0) step.classList.add('active');
            });
        }
        
        function initializeEditor(editorNum) {
            const editor = editors[editorNum];
            
            // Initialize canvas
            editor.canvas = document.getElementById(`imageCanvas${editorNum}`);
            editor.ctx = editor.canvas.getContext('2d');
            
            // Upload area functionality
            const uploadArea = document.getElementById(`uploadArea${editorNum}`);
            const imageUpload = document.getElementById(`imageUpload${editorNum}`);
            
            uploadArea.addEventListener('click', () => imageUpload.click());
            
            // Drag and drop
            uploadArea.addEventListener('dragover', (e) => {
                e.preventDefault();
                uploadArea.classList.add('dragover');
            });
            
            uploadArea.addEventListener('dragleave', () => {
                uploadArea.classList.remove('dragover');
            });
            
            uploadArea.addEventListener('drop', (e) => {
                e.preventDefault();
                uploadArea.classList.remove('dragover');
                const files = e.dataTransfer.files;
                if (files.length > 0) {
                    handleImageUpload(files[0], editorNum);
                }
            });
            
            imageUpload.addEventListener('change', (e) => {
                if (e.target.files.length > 0) {
                    handleImageUpload(e.target.files[0], editorNum);
                }
            });
            
            // Zoom control
            document.getElementById(`zoomSlider${editorNum}`).addEventListener('input', (e) => {
                editor.zoom = parseFloat(e.target.value);
                document.getElementById(`zoomValue${editorNum}`).textContent = Math.round(editor.zoom * 100) + '%';
                redrawCanvas(editorNum);
                updateCropData(editorNum);
            });
            
            // Canvas drag events
            editor.canvas.addEventListener('mousedown', (e) => startDrag(e, editorNum));
            editor.canvas.addEventListener('mousemove', (e) => updateDrag(e, editorNum));
            editor.canvas.addEventListener('mouseup', (e) => endDrag(e, editorNum));
            editor.canvas.addEventListener('mouseleave', (e) => endDrag(e, editorNum));
            
            // Prevent context menu
            editor.canvas.addEventListener('contextmenu', (e) => e.preventDefault());
        }
        
        function handleImageUpload(file, editorNum) {
            if (!file.type.startsWith('image/')) {
                alert(`Please upload a valid image file for Image ${editorNum}.`);
                return;
            }
            
            // Check file size (50MB limit for safety)
            if (file.size > 50 * 1024 * 1024) {
                alert(`Image ${editorNum} is too large. Please use an image smaller than 50MB.`);
                return;
            }
            
            const editor = editors[editorNum];
            
            // Show progress
            showProgress(0, `Reading image ${editorNum} file...`, editorNum);
            
            const reader = new FileReader();
            reader.onload = function(e) {
                showProgress(30, `Processing image ${editorNum}...`, editorNum);
                
                const img = new Image();
                img.onload = function() {
                    showProgress(60, `Optimizing image ${editorNum} for web editor...`, editorNum);
                    
                    setTimeout(() => {
                        // If image is very large, create a compressed version for editor
                        if (img.width > 2000 || img.height > 2000) {
                            showProgress(80, `Compressing large image ${editorNum}...`, editorNum);
                            
                            const compressionCanvas = document.createElement('canvas');
                            const compressionCtx = compressionCanvas.getContext('2d');
                            
                            // Calculate new size (max 2000px on longest side for editor performance)
                            const maxSize = 2000;
                            let newWidth = img.width;
                            let newHeight = img.height;
                            
                            if (img.width > img.height) {
                                if (img.width > maxSize) {
                                    newWidth = maxSize;
                                    newHeight = (img.height * maxSize) / img.width;
                                }
                            } else {
                                if (img.height > maxSize) {
                                    newHeight = maxSize;
                                    newWidth = (img.width * maxSize) / img.height;
                                }
                            }
                            
                            compressionCanvas.width = newWidth;
                            compressionCanvas.height = newHeight;
                            compressionCtx.drawImage(img, 0, 0, newWidth, newHeight);
                            
                            // Create new image from compressed canvas
                            const compressedImage = new Image();
                            compressedImage.onload = function() {
                                editor.originalImage = compressedImage;
                                editor.currentImage = compressedImage;
                                showProgress(100, 'Ready!', editorNum);
                                setTimeout(() => {
                                    hideProgress(editorNum);
                                    showEditor(editorNum);
                                }, 500);
                            };
                            compressedImage.src = compressionCanvas.toDataURL('image/jpeg', 0.85);
                        } else {
                            editor.originalImage = img;
                            editor.currentImage = img;
                            showProgress(100, 'Ready!', editorNum);
                            setTimeout(() => {
                                hideProgress(editorNum);
                                showEditor(editorNum);
                            }, 300);
                        }
                    }, 100);
                };
                img.src = e.target.result;
            };
            reader.readAsDataURL(file);
        }
        
        function showProgress(percentage, text, editorNum) {
            document.getElementById(`uploadProgress${editorNum}`).style.display = 'block';
            document.getElementById(`progressBar${editorNum}`).style.width = percentage + '%';
            document.getElementById(`progressText${editorNum}`).textContent = text;
        }
        
        function hideProgress(editorNum) {
            document.getElementById(`uploadProgress${editorNum}`).style.display = 'none';
        }
        
        function showEditor(editorNum) {
            // Show editor
            document.getElementById(`imageEditor${editorNum}`).style.display = 'block';
            // Update section to show it has image
            document.getElementById(`imageSection${editorNum}`).classList.add('has-image');
            
            // Auto-fit the image initially
            autoFit(editorNum);
            
            redrawCanvas(editorNum);
            updateCropData(editorNum);
        }
        
        function redrawCanvas(editorNum) {
            const editor = editors[editorNum];
            if (!editor.currentImage) return;
            
            // Clear canvas
            editor.ctx.clearRect(0, 0, editor.canvas.width, editor.canvas.height);
            
            // Draw background pattern
            editor.ctx.fillStyle = '#f8f9fa';
            editor.ctx.fillRect(0, 0, editor.canvas.width, editor.canvas.height);
            
            // Calculate image display size
            const displayWidth = editor.currentImage.width * editor.zoom;
            const displayHeight = editor.currentImage.height * editor.zoom;
            
            // Calculate image position (centered + offset from dragging)
            const imageX = (editor.canvas.width - displayWidth) / 2 + editor.offsetX;
            const imageY = (editor.canvas.height - displayHeight) / 2 + editor.offsetY;
            
            // Draw the image
            editor.ctx.drawImage(editor.currentImage, imageX, imageY, displayWidth, displayHeight);
            
            // Calculate the EXACT frame position (matching crop calculation)
            const frameX = (editor.canvas.width - editor.FRAME_WIDTH) / 2;
            const frameY = (editor.canvas.height - editor.FRAME_HEIGHT) / 2;
            
            // Draw frame overlay - Frame showing exact crop area with brand color
            const frameColor = '#F454CD';
            editor.ctx.strokeStyle = frameColor;
            editor.ctx.lineWidth = 4;
            editor.ctx.setLineDash([]);
            editor.ctx.strokeRect(frameX, frameY, editor.FRAME_WIDTH, editor.FRAME_HEIGHT);
            
            // Add inner border for better visibility
            editor.ctx.strokeStyle = 'rgba(255, 255, 255, 0.8)';
            editor.ctx.lineWidth = 2;
            editor.ctx.strokeRect(frameX + 2, frameY + 2, editor.FRAME_WIDTH - 4, editor.FRAME_HEIGHT - 4);
            
            // Add frame label
            const labelText = editorNum === 1 ? '19.05 cm × 10.79 cm' : '17.69 cm × 11.38 cm';
            const labelWidth = 140;
            editor.ctx.fillStyle = frameColor;
            editor.ctx.fillRect(frameX, frameY - 25, labelWidth, 20);
            editor.ctx.fillStyle = 'white';
            editor.ctx.font = 'bold 12px Inter, sans-serif';
            editor.ctx.fillText(labelText, frameX + 5, frameY - 10);
        }
        
        function startDrag(e, editorNum) {
            const editor = editors[editorNum];
            if (!editor.currentImage) return;
            
            const rect = editor.canvas.getBoundingClientRect();
            editor.lastMouseX = e.clientX - rect.left;
            editor.lastMouseY = e.clientY - rect.top;
            editor.isDragging = true;
            editor.canvas.style.cursor = 'grabbing';
        }
        
        function updateDrag(e, editorNum) {
            const editor = editors[editorNum];
            if (!editor.isDragging || !editor.currentImage) return;
            
            const rect = editor.canvas.getBoundingClientRect();
            const mouseX = e.clientX - rect.left;
            const mouseY = e.clientY - rect.top;
            
            // Calculate drag distance
            const deltaX = mouseX - editor.lastMouseX;
            const deltaY = mouseY - editor.lastMouseY;
            
            // Update image position
            editor.offsetX += deltaX;
            editor.offsetY += deltaY;
            
            // Update last mouse position
            editor.lastMouseX = mouseX;
            editor.lastMouseY = mouseY;
            
            redrawCanvas(editorNum);
            updateCropData(editorNum);
        }
        
        function endDrag(e, editorNum) {
            const editor = editors[editorNum];
            editor.isDragging = false;
            editor.canvas.style.cursor = 'grab';
        }
        
        function updateCropData(editorNum) {
            const editor = editors[editorNum];
            if (!editor.currentImage) return;
            
            // Calculate frame position (EXACTLY matching the visual frame)
            const frameX = (editor.canvas.width - editor.FRAME_WIDTH) / 2;
            const frameY = (editor.canvas.height - editor.FRAME_HEIGHT) / 2;
            
            // Calculate image position (EXACTLY matching the visual image)
            const displayWidth = editor.currentImage.width * editor.zoom;
            const displayHeight = editor.currentImage.height * editor.zoom;
            const imageX = (editor.canvas.width - displayWidth) / 2 + editor.offsetX;
            const imageY = (editor.canvas.height - displayHeight) / 2 + editor.offsetY;
            
            // Calculate crop coordinates in original image space
            let cropX = (frameX - imageX) / editor.zoom;
            let cropY = (frameY - imageY) / editor.zoom;
            let cropWidth = editor.FRAME_WIDTH / editor.zoom;
            let cropHeight = editor.FRAME_HEIGHT / editor.zoom;
            
            // Ensure crop stays within image bounds
            cropX = Math.max(0, Math.min(cropX, editor.currentImage.width - cropWidth));
            cropY = Math.max(0, Math.min(cropY, editor.currentImage.height - cropHeight));
            cropWidth = Math.min(cropWidth, editor.currentImage.width - cropX);
            cropHeight = Math.min(cropHeight, editor.currentImage.height - cropY);
            
            const cropData = {
                x: cropX,
                y: cropY,
                width: cropWidth,
                height: cropHeight,
                frameX: frameX,
                frameY: frameY,
                imageX: imageX,
                imageY: imageY,
                zoom: editor.zoom
            };
            
            document.getElementById(`cropCoordinates${editorNum}`).value = JSON.stringify(cropData);
            
            // Generate cropped image data with exact coordinates
            generateCroppedImageData(cropData, editorNum);
            
            // Update debug display
            document.getElementById(`cropCoords${editorNum}`).textContent = 
                `${cropX.toFixed(1)}, ${cropY.toFixed(1)}, ${cropWidth.toFixed(1)}×${cropHeight.toFixed(1)}`;
        }
        
        function generateCroppedImageData(cropData, editorNum) {
            const editor = editors[editorNum];
            const croppedCanvas = document.createElement('canvas');
            const croppedCtx = croppedCanvas.getContext('2d');
            
            // Set canvas to target dimensions - smaller for upload efficiency
            const targetWidth = 600;  // Smaller for upload, still good quality
            const targetHeight = Math.round(targetWidth / editor.TARGET_RATIO);
            
            croppedCanvas.width = targetWidth;
            croppedCanvas.height = targetHeight;
            
            // Draw the cropped portion
            croppedCtx.drawImage(
                editor.currentImage,
                cropData.x, cropData.y, cropData.width, cropData.height,
                0, 0, targetWidth, targetHeight
            );
            
            // Aggressive compression to prevent 413 errors
            let quality = 0.6;  // Start with lower quality
            let imageData = croppedCanvas.toDataURL('image/jpeg', quality);
            
            // Keep reducing until under 1MB
            while (imageData.length > 1 * 1024 * 1024 && quality > 0.2) {
                quality -= 0.1;
                imageData = croppedCanvas.toDataURL('image/jpeg', quality);
            }
            
            // If still too large, reduce dimensions
            if (imageData.length > 1 * 1024 * 1024) {
                const smallerCanvas = document.createElement('canvas');
                const smallerCtx = smallerCanvas.getContext('2d');
                smallerCanvas.width = 400;
                smallerCanvas.height = Math.round(400 / editor.TARGET_RATIO);
                
                smallerCtx.drawImage(croppedCanvas, 0, 0, smallerCanvas.width, smallerCanvas.height);
                imageData = smallerCanvas.toDataURL('image/jpeg', 0.7);
            }
            
            console.log(`Image ${editorNum} final upload size: ${Math.round(imageData.length / 1024)}KB`);
            
            // Store the cropped image data
            if (editorNum === 1) {
                document.getElementById('croppedImageData1').value = imageData;
            } else {
                document.getElementById('croppedImageData2').value = imageData;
            }
        }
        
        function resetPosition(editorNum) {
            const editor = editors[editorNum];
            editor.zoom = 1;
            editor.offsetX = 0;
            editor.offsetY = 0;
            document.getElementById(`zoomSlider${editorNum}`).value = 1;
            document.getElementById(`zoomValue${editorNum}`).textContent = '100%';
            redrawCanvas(editorNum);
            updateCropData(editorNum);
        }
        
        function autoFit(editorNum) {
            const editor = editors[editorNum];
            if (!editor.currentImage) return;
            
            // Calculate optimal zoom to fit image in frame with some padding
            const frameRatio = editor.FRAME_WIDTH / editor.FRAME_HEIGHT;
            const imageRatio = editor.currentImage.width / editor.currentImage.height;
            
            if (imageRatio > frameRatio) {
                // Image is wider relative to frame - fit to frame width
                editor.zoom = (editor.FRAME_WIDTH * 1.2) / editor.currentImage.width;  // 1.2 for some padding
            } else {
                // Image is taller relative to frame - fit to frame height
                editor.zoom = (editor.FRAME_HEIGHT * 1.2) / editor.currentImage.height;
            }
            
            // Constrain zoom to reasonable bounds
            editor.zoom = Math.max(0.1, Math.min(3, editor.zoom));
            
            // Center the image
            editor.offsetX = 0;
            editor.offsetY = 0;
            
            // Update UI
            document.getElementById(`zoomSlider${editorNum}`).value = editor.zoom;
            document.getElementById(`zoomValue${editorNum}`).textContent = Math.round(editor.zoom * 100) + '%';
            
            redrawCanvas(editorNum);
            updateCropData(editorNum);
        }
        
        // TNB NA dependency function
        function updateTnbNa() {
            const tnbMeter = document.getElementById('tnb_meter').value;
            const tnbNaField = document.getElementById('tnb_na');
            
            if (tnbMeter === 'TNB Tenant Meter') {
                tnbNaField.value = 'CNG';
            } else if (tnbMeter === 'kWh Meter') {
                tnbNaField.value = 'NA';
            }
        }
        
        // Form validation
        document.getElementById('proposalForm').addEventListener('submit', function(e) {
            // Check if images are uploaded
            const imageData1 = document.getElementById('croppedImageData1').value;
            const imageData2 = document.getElementById('croppedImageData2').value;
            
            if (!imageData1 && !imageData2) {
                const proceed = confirm('No images uploaded. Do you want to generate the proposal without any images?');
                if (!proceed) {
                    e.preventDefault();
                    return;
                }
            }
            
            // Disable submit button during processing
            const submitBtn = document.getElementById('generateBtn');
            submitBtn.disabled = true;
            submitBtn.textContent = 'Generating Proposal...';
        });
    </script>
</body>
</html>
    '''

@app.route('/generate', methods=['POST'])
def generate_proposal():
    """
    Generate the proposal from form data including both processed images
    """
    try:
        print("📝 Processing form submission with dual images...")
        
        # Get form data
        form_data = {
            'building_name': request.form.get('building_name', '').strip(),
            'type_building': request.form.get('type_building', ''),
            'address': request.form.get('address', '').strip(),
            'survey_date': request.form.get('survey_date', ''),
            'prepared_by': request.form.get('prepared_by', '').strip(),
            'prepared_date': request.form.get('prepared_date', ''),
            'building_manager_name': request.form.get('building_manager_name', '').strip(),
            'building_manager_email': request.form.get('building_manager_email', '').strip(),
            'building_manager_phone': request.form.get('building_manager_phone', '').strip(),
            'building_manager_company': request.form.get('building_manager_company', '').strip(),
            'otic': request.form.get('otic', '').strip(),
            'tap_new_or_spare': request.form.get('tap_new_or_spare', ''),
            'tapping_location': request.form.get('tapping_location', '').strip(),
            'tapping_location_level': request.form.get('tapping_location_level', '').strip(),
            'site_assessment_mccb': request.form.get('site_assessment_mccb', '').strip(),
            'tnb_meter': request.form.get('tnb_meter', ''),
            'tnb_na': request.form.get('tnb_na', '').strip(),
            'parking_location': request.form.get('parking_location', '').strip(),
            'ev_charger_model': request.form.get('ev_charger_model', ''),
            'network_strength': request.form.get('network_strength', '')
        }
        
        print(f"📋 Form data received: {form_data['prepared_by']} - {form_data['address'][:50]}...")
        print(f"🏢 Building Type received: '{form_data['type_building']}'")
        print(f"📊 All form data: {form_data}")
        
        # Format dates
        if form_data['survey_date']:
            survey_date_obj = datetime.strptime(form_data['survey_date'], '%Y-%m-%d')
            form_data['survey_date'] = survey_date_obj.strftime('%B %d, %Y')
        
        if form_data['prepared_date']:
            prepared_date_obj = datetime.strptime(form_data['prepared_date'], '%Y-%m-%d')
            form_data['prepared_date'] = prepared_date_obj.strftime('%B %d, %Y')
        
        # Process both uploaded images
        image_path = None
        image_path_2 = None
        
        # Process first image
        cropped_image_data = request.form.get('cropped_image_data')
        crop_coordinates = request.form.get('crop_coordinates')
        
        if cropped_image_data and crop_coordinates:
            print("🖼️  Processing first uploaded image...")
            try:
                crop_data = json.loads(crop_coordinates)
                image_path, img_error = process_cropped_image(cropped_image_data, crop_data, "1")
                
                if img_error:
                    print(f"❌ First image processing error: {img_error}")
                    flash(f"First image processing error: {img_error}", 'error')
                    return redirect(url_for('index'))
                else:
                    print(f"✅ First image processed successfully: {image_path}")
                    
            except json.JSONDecodeError as json_err:
                print(f"❌ Invalid first image coordinates: {json_err}")
                flash("Invalid first image coordinates. Please try uploading the image again.", 'error')
                return redirect(url_for('index'))
            except Exception as img_e:
                print(f"❌ First image processing failed: {str(img_e)}")
                flash(f"First image processing failed: {str(img_e)}", 'error')
                return redirect(url_for('index'))
        
        # Process second image
        cropped_image_data_2 = request.form.get('cropped_image_data_2')
        crop_coordinates_2 = request.form.get('crop_coordinates_2')
        
        if cropped_image_data_2 and crop_coordinates_2:
            print("🖼️  Processing second uploaded image...")
            try:
                crop_data_2 = json.loads(crop_coordinates_2)
                image_path_2, img_error_2 = process_cropped_image(cropped_image_data_2, crop_data_2, "2")
                
                if img_error_2:
                    print(f"❌ Second image processing error: {img_error_2}")
                    flash(f"Second image processing error: {img_error_2}", 'error')
                    return redirect(url_for('index'))
                else:
                    print(f"✅ Second image processed successfully: {image_path_2}")
                    
            except json.JSONDecodeError as json_err:
                print(f"❌ Invalid second image coordinates: {json_err}")
                flash("Invalid second image coordinates. Please try uploading the image again.", 'error')
                return redirect(url_for('index'))
            except Exception as img_e:
                print(f"❌ Second image processing failed: {str(img_e)}")
                flash(f"Second image processing failed: {str(img_e)}", 'error')
                return redirect(url_for('index'))
        
        if not image_path and not image_path_2:
            print("ℹ️  No images provided - generating text-only proposal")
        
        # Generate output filename
        client_name = form_data['address'].split(',')[0].strip()[:20]
        safe_client_name = "".join(c for c in client_name if c.isalnum() or c in (' ', '_')).strip()
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f"proposal_{safe_client_name.replace(' ', '_')}_{timestamp}.pptx"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        
        print(f"📄 Generating proposal with dual images: {output_filename}")
        
        # Generate the proposal (now with both image paths)
        success, message = replace_placeholders_and_images_in_pptx(
            TEMPLATE_PATH, form_data, image_path, image_path_2, output_path
        )
        
        # Clean up temporary images
        for temp_image in [image_path, image_path_2]:
            if temp_image and os.path.exists(temp_image):
                try:
                    os.remove(temp_image)
                    print("🗑️  Cleaned up temporary image file")
                except:
                    pass  # Ignore cleanup errors
        
        if success:
            print(f"✅ Proposal generated successfully: {output_path}")
            return send_file(output_path, as_attachment=True, download_name=output_filename)
        else:
            print(f"❌ Generation failed: {message}")
            flash(f"Error: {message}", 'error')
            return redirect(url_for('index'))
            
    except Exception as e:
        print(f"❌ Unexpected error: {str(e)}")
        import traceback
        traceback.print_exc()
        flash(f"Unexpected error: {str(e)}", 'error')
        return redirect(url_for('index'))

@app.errorhandler(413)
def too_large(e):
    """
    Handle file too large errors
    """
    print("❌ File too large error caught")
    flash("Image file is too large. Please use a smaller image (under 50MB).", 'error')
    return redirect(url_for('index'))

@app.route('/shutdown')
def shutdown():
    """
    Shutdown the server
    """
    print("🔥 Shutting down server...")
    func = request.environ.get('werkzeug.server.shutdown')
    if func is None:
        raise RuntimeError('Not running with the Werkzeug Server')
    func()
    return 'Server shutting down...'

def open_browser():
    """
    Open browser after short delay
    """
    time.sleep(1.5)
    webbrowser.open('http://localhost:5000')

def main():
    """
    Main function to start the application
    """
    print("🚀 Starting Dual Image Proposal Generator...")
    print("📁 Template path:", TEMPLATE_PATH)
    
    # Check if template exists
    if not os.path.exists(TEMPLATE_PATH):
        print(f"❌ ERROR: Template file not found at {TEMPLATE_PATH}")
        print("📍 Please make sure the template file exists at the specified location.")
        input("Press Enter to exit...")
        return
    
    print("✅ Template file found!")
    print("🌐 Starting web server...")
    print("📂 Generated proposals will be saved to:", os.path.abspath(OUTPUT_FOLDER))
    print("🖼️  Temporary images will be saved to:", os.path.abspath(TEMP_IMAGES_FOLDER))
    print("\n" + "="*60)
    print("🎯 READY! Browser will open automatically...")
    print("📝 Fill out the form")
    print("🖼️  Upload and crop your images for both placeholders") 
    print("   - IMG_PLACEHOLDER: 19.05 cm × 10.79 cm")
    print("   - IMG_PLACEHOLDER2: 17.69 cm × 11.38 cm")
    print("🚀 Generate your proposal")
    print("🔗 If browser doesn't open, go to: http://localhost:5000")
    print("🛑 To stop: use Ctrl+C or close the browser tab")
    print("="*60 + "\n")
    
    # Start browser in separate thread
    #threading.Thread(target=open_browser, daemon=True).start()
    
    # Start Flask app
    try:
        app.run(host='localhost', port=5000, debug=False, use_reloader=False)
    except KeyboardInterrupt:
        print("\n🔥 Application stopped by user")
    except Exception as e:
        print(f"❌ Error starting server: {e}")
        input("Press Enter to exit...")

if __name__ == '__main__':
    main()